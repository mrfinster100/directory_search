## directory_search.py
## Searches files in directory and subdirectories for a specified string.

## Setup:
# 1) Run file.reg to enable drag-and-drop conversion
# 2) Install necessary Python 3.8 dependencies: (i.e. PyPDF2, docx, xlrd, python-pptx)
# 3) Drag and drop any parent directory or file on top of directory_search.py
# 4) Specify a searchable string (not case-sensitive)
# 5) View list of files containing the specified string w/ instances
##TODO:
#Auto adjust width and height of cells
#replace commas (in context list) with line breaks for xlsx output
# if pdf requires decryption, skip
# handling of protected mode files
# handling of unsupported file types


from docx import Document
from pptx import Presentation
import sys
import os
import re
import PyPDF2
import csv
import xlrd
import xml.etree.ElementTree as ET
import ctypes
import pandas as pd

## 1
# 3rd Level Helper Function
# for searching ppt/pptx
def search_ppt(pptfname):
      array = []
      prs = Presentation(pptfname)
      for slide in prs.slides:
            for shape in slide.shapes:
                  if not shape.has_text_frame:
                      continue
                  for paragraph in shape.text_frame.paragraphs:
                        #for run in paragraph.runs:
                        paratext = paragraph.text
                        words = paratext.split()
                        if specifiedString in words:
                              array.append(paratext)
                        
      text = (" ").join(array)
      matches = re.findall(specifiedString, text, flags=re.IGNORECASE)      
      instances = len(matches)
     
      csvdatapoints = [array, instances]
      return csvdatapoints

## 2
# 3rd Level Helper Function
# for searching doc/docx
def search_doc(docfname):
      array = []
      document = Document(docfname)
      NumParagraphs = len(document.paragraphs)
      text = ""
      for i in range(0, NumParagraphs):
            currentPara = document.paragraphs[i].text
            if specifiedString in currentPara:
                  array.append(currentPara)
      text = ("\n").join(array)
      matches = re.findall(specifiedString, text, flags=re.IGNORECASE)
      instances = len(matches)
      
      csvdatapoints = [array, instances]
      return csvdatapoints
    
## 3
# 3rd Level Helper Function
# for searching .pdf files##DOESN'T WORK FOR ALL PDF FILES!
def search_pdf(pdffname):
      text = ""
      array = []
      pdfFileObj = open(pdffname, 'rb')
      pdfReader = PyPDF2.PdfFileReader(pdfFileObj)
      NumPages = pdfReader.getNumPages()
      for i in range(0, NumPages):
            PageObj = pdfReader.getPage(i)
            pgtxt = PageObj.extractText()
            text += pgtxt
            print(pgtxt)
            if specifiedString in pgtxt:
                  array.append(pgtxt)
      text = text.split()
      text = (" ").join(text)
      matches = re.findall(specifiedString, text, flags=re.IGNORECASE)
      instances = len(matches)
      
      csvdatapoints = [array, instances]
      return csvdatapoints

      

## 4  
# 3rd Level Helper Function    
# for searching .csv files
def search_csv(csvfname):
      text = ""
      array = []
      with open(csvfname, newline='') as f:
            reader = csv.reader(f, delimiter=',',quotechar="|")
            for row in reader:
                  for column in row:
                        if specifiedString in column:
                              array.append(column)
                              text+=column
      matches = re.findall(specifiedString, text, flags=re.IGNORECASE)
      instances = len(matches)
      
      csvdatapoints = [array, instances]
      return csvdatapoints
            
## 5         
# 3rd Level Helper Function
# for searching .xlsx
def search_xlsx(xlsxfname):
      array = []
      xl_workbook = xlrd.open_workbook(xlsxfname, on_demand=True)
      res = len(xl_workbook.sheet_names()) 
      text = ""
      for i in range(0, res):
            xl_sheet = xl_workbook.sheet_by_index(i)
            num_cols = xl_sheet.ncols
            for row_idx in range(0, xl_sheet.nrows):
                  for col_idx in range(0, num_cols): 
                        cell_obj = xl_sheet.cell(row_idx, col_idx)
                        converted2string = str(cell_obj.value)
                      
                        if specifiedString in converted2string:
                              array.append(converted2string)
                              
      text = (" ").join(array)
      matches = re.findall(specifiedString, text, flags=re.IGNORECASE)
      instances = len(matches)
      csvdatapoints = [array, instances]
      return csvdatapoints
      
## 1
# 2nd Level Switch Function   
# receives a file directory
# returns # of matches
def fileSwitch(fileDir): 
    pre, ext = os.path.splitext(fileDir)
    # search pptx/ppt
    if ext == '.pptx' or ext == '.ppt':
        return search_ppt(fileDir)
    # search doc/docx
    if ext == ".doc" or ext == '.docx':
        return search_doc(fileDir)
    # search pdf
    if ext == ".pdf":
        return search_pdf(fileDir)
    # search csv
    if ext == ".csv":
        return search_csv(fileDir)
    ## search xlsx
    if ext == ".xlsx":
        return search_xlsx(fileDir)
## 1
##MAIN FUNCTION##
      
data = {}
sentenceWithHighlights = []
instances = []
filenames = []
try:
      if len(sys.argv) != 1:
            specifiedDirectory = sys.argv[1]#"C:\Python_Scripts\directory_search-master\clean"
            specifiedString = input("What string would you like to search for?\n")
            for subdir, dirs, files in os.walk(specifiedDirectory):
                  for file in files:
                        file = os.path.join(subdir, file)
                        datapoint = fileSwitch(file)##Returns  [sentencesWithHighlights, instances]
                        smh = datapoint[0]
                        instc = datapoint[1]
                        relPath = os.path.relpath(file, subdir)
                        filenames.append(relPath)
                        instances.append(str(instc))
                        sentenceWithHighlights.append(smh)
                        

            #push datapoints to pandas dataframe
            data["Found In"]=filenames
            data["Frequency"]=instances
            data["Context"]=sentenceWithHighlights
            
            
            df = pd.DataFrame(data)

            #create xlsx file with pandas dataframe
            writer = pd.ExcelWriter(specifiedDirectory+'_search'+'.xlsx',engine='xlsxwriter')#add timestamp
            df.to_excel(writer,sheet_name='Sheet1')
            writer.save()
         
            print("Job completed.")
            input("Press any key to exit...")
      else:
          ctypes.windll.user32.MessageBoxW(None, "Drag a folder or file onto directory_search.py \nto search for a specified string.", "Ok", 0)          
          exit
except Exception as ex:
    print(ex)
    input()
